from __future__ import annotations

from typing import List, Dict, Any
import os

from pptx import Presentation
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.shapes.base import BaseShape

from ..core.Models import Document


def _extract_slide_text(slide) -> str:
    """
    Extract only real, visible text from a PPTX slide.
    Filters out:
    - empty shapes
    - shapes without .text attribute
    - shapes that only contain filenames or XML placeholders
    """

    texts: List[str] = []

    for shape in slide.shapes:

        # Skip pictures, charts, tables, etc.
        if not hasattr(shape, "text_frame"):
            # Some shapes have .text attribute but it's garbage XML, skip them
            if hasattr(shape, "text") and isinstance(shape.text, str):
                t = shape.text.strip()
                # Only keep real text, not filenames or slide notes
                if len(t) > 0 and not t.endswith(".pptx"):
                    texts.append(t)
            continue

        # A proper text frame
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)

    # Merge all collected text lines
    cleaned = "\n".join(texts).strip()

    return cleaned


def load_pptx_as_documents(
    path: str,
    per_slide: bool = True,
    extra_metadata: Dict[str, Any] | None = None,
) -> List[Document]:
    """
    Load a PPTX file and return clean Document objects.
    """

    if not os.path.isfile(path):
        raise FileNotFoundError(f"PPTX file not found: {path}")

    extra_metadata = extra_metadata or {}
    prs = Presentation(path)

    documents: List[Document] = []

    if per_slide:
        for i, slide in enumerate(prs.slides):
            text = _extract_slide_text(slide)
            if not text:
                continue

            doc_id = f"{os.path.basename(path)}_slide_{i}"
            metadata = {
                "source": os.path.abspath(path),
                "slide": i,
                **extra_metadata,
            }

            documents.append(
                Document(
                    id=doc_id,
                    text=text,
                    metadata=metadata,
                )
            )
    else:
        all_text = []
        for slide in prs.slides:
            t = _extract_slide_text(slide)
            if t:
                all_text.append(t)

        full = "\n\n".join(all_text).strip()
        if full:
            doc_id = os.path.basename(path)
            metadata = {
                "source": os.path.abspath(path),
                **extra_metadata,
            }
            documents.append(
                Document(
                    id=doc_id,
                    text=full,
                    metadata=metadata,
                )
            )

    return documents
